from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def add_title_slide(prs, title, subtitle=None):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if subtitle is not None:
        slide.placeholders[1].text = subtitle


def add_bullets_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.level = 0
    return slide


def add_two_column_slide(prs, title, left_bullets, right_bullets):
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    # Left textbox
    left = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(4.3), Inches(4.5))
    ltf = left.text_frame
    ltf.clear()
    for i, t in enumerate(left_bullets):
        p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
        p.text = t
        p.level = 0

    # Right textbox
    right = slide.shapes.add_textbox(Inches(5.1), Inches(1.6), Inches(4.3), Inches(4.5))
    rtf = right.text_frame
    rtf.clear()
    for i, t in enumerate(right_bullets):
        p = rtf.paragraphs[0] if i == 0 else rtf.add_paragraph()
        p.text = t
        p.level = 0

    return slide


def build_presentation(output_path: str = "Uniform_Detector_Review2.pptx"):
    prs = Presentation()

    # 1) Title
    add_title_slide(
        prs,
        title="Student Uniform Detection – Review 2",
        subtitle="YOLO-based detection | Laptop webcam | Arduino flag output",
    )

    # 2) System Architecture
    add_bullets_slide(
        prs,
        title="System Architecture",
        bullets=[
            "Input: Laptop Webcam / Video File",
            "Model Inference: YOLOv11 (Ultralytics) with fallbacks (yolo11n/yolov8m)",
            "Detection Normalization: Maps labels to standard set",
            "Rule Engine: Boys (ID+Shirt+Pant+Shoes), Girls (ID+Top+Pant+Shoes)",
            "Output: On-frame overlay + Terminal status (1/0) + Serial to Arduino",
            "Artifacts: Trained weights in runs/train/.../best.pt",
        ],
    )

    # 3) Complete Workflow
    add_bullets_slide(
        prs,
        title="Complete Workflow",
        bullets=[
            "Start app → Resolve model path relative to script",
            "Open camera → Read frames (every 5th frame for inference)",
            "Run YOLO → Collect class names + confidences",
            "Normalize detections → Aggregate counts",
            "Apply uniform rules → Decide type (Boys/Girls) + completeness",
            "Report: draw message, print flag (1/0), send over serial if enabled",
        ],
    )

    # 4) Additional Work Since Review 1
    add_bullets_slide(
        prs,
        title="Additional Work Since Review 1",
        bullets=[
            "Project restructured: essential files in root, archives in scrap_files/",
            "Run scripts fixed (no hardcoded subfolders)",
            "Robust model path resolution + ordered fallbacks",
            "Detection count logic refined to prevent double-counting",
            "Webcam handling guarded with try/finally to release resources",
            "Serial support enabled via pyserial; auto-detect Arduino ports",
            "Documentation improved; repo published to GitHub",
        ],
    )

    # 5) Results / Demo Notes
    add_two_column_slide(
        prs,
        title="Results & Demo Notes",
        left_bullets=[
            "Live detection via laptop webcam",
            "Status overlay and terminal flag printed",
            "Works with trained weights (preferred) or fallback models",
        ],
        right_bullets=[
            "Classes: Identity Card, Shirt/Top, Pant, Shoes, Slippers",
            "Threshold: 0.5 (adjustable)",
            "Flag: 1 for complete, 0 otherwise; sent to Arduino",
        ],
    )

    # 6) Next Steps (optional but fits 6-slide max)
    add_bullets_slide(
        prs,
        title="Next Steps",
        bullets=[
            "Broaden dataset and retrain for higher accuracy",
            "Add simple UI toggle for camera/device selection",
            "Export lightweight build for edge devices (e.g., Jetson/RPi)",
            "Optional: web API for remote triggering and monitoring",
        ],
    )

    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    path = build_presentation()
    print(f"✅ Generated: {path}")
